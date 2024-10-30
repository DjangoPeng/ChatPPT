import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from utils import remove_all_slides
from logger import LOG  # 引入日志模块

def format_text(paragraph, text):
    """
    格式化文本，处理加粗内容。假设 ** 包围的文本表示需要加粗。
    """
    while '**' in text:
        start = text.find('**')
        end = text.find('**', start + 2)
        
        if start != -1 and end != -1:
            # 添加加粗之前的普通文本
            if start > 0:
                run = paragraph.add_run()
                run.text = text[:start]
            
            # 添加加粗文本
            bold_run = paragraph.add_run()
            bold_run.text = text[start + 2:end]
            bold_run.font.bold = True  # 设置加粗
            
            # 处理剩余文本
            text = text[end + 2:]
        else:
            break

    # 添加剩余的普通文本
    if text:
        run = paragraph.add_run()
        run.text = text

def insert_image_centered_in_placeholder(new_slide, image_path):
    """
    将图片插入到 Slide 中，使其中心与 placeholder 的中心对齐。
    如果图片尺寸超过 placeholder，则进行缩小适配。
    在插入成功后删除 placeholder。
    """
    # 构建图片的绝对路径
    image_full_path = os.path.join(os.getcwd(), image_path)
    
    # 检查图片是否存在
    if not os.path.exists(image_full_path):
        LOG.warning(f"图片路径 '{image_full_path}' 不存在，跳过此图片。")
        return

    # 打开图片并获取其大小（以像素为单位）
    with Image.open(image_full_path) as img:
        img_width_px, img_height_px = img.size

    # 遍历找到图片的 placeholder（type 18 表示图片 placeholder）
    for shape in new_slide.placeholders:
        if shape.placeholder_format.type == 18:
            placeholder_width = shape.width
            placeholder_height = shape.height
            placeholder_left = shape.left
            placeholder_top = shape.top

            # 计算 placeholder 的中心点
            placeholder_center_x = placeholder_left + placeholder_width / 2
            placeholder_center_y = placeholder_top + placeholder_height / 2

            # 图片的宽度和高度转换为 PowerPoint 的单位 (Inches)
            img_width = Inches(img_width_px / 96)  # 假设图片 DPI 为 96
            img_height = Inches(img_height_px / 96)

            # 如果图片的宽度或高度超过 placeholder，按比例缩放图片
            if img_width > placeholder_width or img_height > placeholder_height:
                scale = min(placeholder_width / img_width, placeholder_height / img_height)
                img_width *= scale
                img_height *= scale

            # 计算图片左上角位置，使其中心对准 placeholder 中心
            left = placeholder_center_x - img_width / 2
            top = placeholder_center_y - img_height / 2

            # 插入图片到指定位置并设定缩放后的大小
            new_slide.shapes.add_picture(image_full_path, left, top, width=img_width, height=img_height)
            LOG.debug(f"图片已插入，并以 placeholder 中心对齐，路径: {image_full_path}")

            # 移除占位符
            sp = shape._element  # 获取占位符的 XML 元素
            sp.getparent().remove(sp)  # 从父元素中删除
            LOG.debug("已删除图片的 placeholder")
            break

# 生成 PowerPoint 演示文稿
def generate_presentation(powerpoint_data, template_path: str, output_path: str):
    # 检查模板文件是否存在
    if not os.path.exists(template_path):
        LOG.error(f"模板文件 '{template_path}' 不存在。")  # 记录错误日志
        raise FileNotFoundError(f"模板文件 '{template_path}' 不存在。")

    prs = Presentation(template_path)  # 加载 PowerPoint 模板
    remove_all_slides(prs)  # 清除模板中的所有幻灯片
    prs.core_properties.title = powerpoint_data.title  # 设置 PowerPoint 的核心标题

    # 遍历所有幻灯片数据，生成对应的 PowerPoint 幻灯片
    for slide in powerpoint_data.slides:
        # 确保布局索引不超出范围，超出则使用默认布局
        if slide.layout_id >= len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[slide.layout_id]

        new_slide = prs.slides.add_slide(slide_layout)  # 添加新的幻灯片

        # 设置幻灯片标题
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide.content.title
            LOG.debug(f"设置幻灯片标题: {slide.content.title}")

        # 添加文本内容
        for shape in new_slide.shapes:
            # 只处理非标题的文本框
            if shape.has_text_frame and not shape == new_slide.shapes.title:
                text_frame = shape.text_frame
                text_frame.clear()  # 清除原有内容

                # 直接使用第一个段落，不添加新的段落，避免额外空行
                first_paragraph = text_frame.paragraphs[0]
                
                # 将要点内容作为项目符号列表添加到文本框中
                for point in slide.content.bullet_points:
                    # 第一个要点覆盖初始段落，其他要点添加新段落
                    paragraph = first_paragraph if point == slide.content.bullet_points[0] else text_frame.add_paragraph()
                    paragraph.level = point["level"]  # 设置项目符号的级别
                    format_text(paragraph, point["text"])  # 调用 format_text 方法来处理加粗文本
                    LOG.debug(f"添加列表项: {paragraph.text}，级别: {paragraph.level}")

                break

        # 插入图片
        if slide.content.image_path:
            insert_image_centered_in_placeholder(new_slide, slide.content.image_path)

    # 保存生成的 PowerPoint 文件
    prs.save(output_path)
    LOG.info(f"演示文稿已保存到 '{output_path}'")
