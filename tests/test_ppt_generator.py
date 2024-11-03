import unittest
import os
import sys
from pptx import Presentation

# 添加 src 目录到模块搜索路径，以便可以导入 src 目录中的模块
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../src')))

from data_structures import PowerPoint, Slide, SlideContent
from ppt_generator import generate_presentation

class TestPPTGenerator(unittest.TestCase):
    """
    测试 ppt_generator 模块的 generate_presentation 函数，验证生成的 PowerPoint 文件内容是否符合预期。
    """

    def setUp(self):
        """
        设置测试数据和输出路径。
        """
        # 定义输入 PowerPoint 数据结构
        self.powerpoint_data = PowerPoint(
            title="ChatPPT Demo",
            slides=[
                Slide(
                    layout_id=1,
                    layout_name="Title 1",
                    content=SlideContent(title="ChatPPT Demo")
                ),
                Slide(
                    layout_id=2,
                    layout_name="Title, Content 0",
                    content=SlideContent(
                        title="2024 业绩概述",
                        bullet_points=[
                            {"text": "总收入增长15%", "level": 0},
                            {"text": "市场份额扩大至30%", "level": 0}
                        ]
                    )
                ),
                Slide(
                    layout_id=8,
                    layout_name="Title, Content, Picture 2",
                    content=SlideContent(
                        title="业绩图表",
                        bullet_points=[{"text": "OpenAI 利润不断增加", "level": 0}],
                        image_path="images/performance_chart.png"
                    )
                ),
                Slide(
                    layout_id=8,
                    layout_name="Title, Content, Picture 2",
                    content=SlideContent(
                        title="新产品发布",
                        bullet_points=[
                            {"text": "产品A: **特色功能介绍**", "level": 0},
                            {"text": "增长潜力巨大", "level": 1},
                            {"text": "新兴市场", "level": 1},
                            {"text": "**非洲**市场", "level": 2},
                            {"text": "**东南亚**市场", "level": 2},
                            {"text": "产品B: 市场定位", "level": 0}
                        ],
                        image_path="images/forecast.png"
                    )
                )
            ]
        )

        self.template_path = "templates/SimpleTemplate.pptx"  # 假设存在模板文件
        self.output_path = "outputs/test_presentation.pptx"  # 定义输出文件路径

    def test_generate_presentation(self):
        """
        测试 generate_presentation 函数生成的 PowerPoint 文件是否符合预期。
        """
        # 调用函数生成 PowerPoint 演示文稿
        generate_presentation(self.powerpoint_data, self.template_path, self.output_path)

        # 检查输出文件是否存在
        self.assertTrue(os.path.exists(self.output_path), "输出 PowerPoint 文件未找到。")

        # 打开生成的 PowerPoint 文件并验证内容
        prs = Presentation(self.output_path)
        
        # 检查演示文稿标题
        self.assertEqual(prs.core_properties.title, self.powerpoint_data.title)

        # 检查幻灯片数量
        self.assertEqual(len(prs.slides), len(self.powerpoint_data.slides))

        # 验证每张幻灯片的内容
        for idx, slide_data in enumerate(self.powerpoint_data.slides):
            slide = prs.slides[idx]

            # 验证幻灯片标题
            self.assertEqual(slide.shapes.title.text, slide_data.content.title)

            # 验证项目符号列表内容
            bullet_points = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame and shape != slide.shapes.title]
            expected_bullets = [point["text"].replace("**", "") for point in slide_data.content.bullet_points]
            for bullet, expected in zip(bullet_points, expected_bullets):
                self.assertIn(expected, bullet)

            # 验证图片路径（如果存在）
            if slide_data.content.image_path:
                images = [shape for shape in slide.shapes if shape.shape_type == 13]  # 13 为图片形状类型
                self.assertGreater(len(images), 0, f"幻灯片 {idx + 1} 应该包含图片，但未找到。")

    def tearDown(self):
        """
        清理生成的文件。
        """
        if os.path.exists(self.output_path):
            os.remove(self.output_path)

if __name__ == "__main__":
    unittest.main()
